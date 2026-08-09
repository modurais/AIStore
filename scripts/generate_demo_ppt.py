from __future__ import annotations

import json
from pathlib import Path
from typing import Dict, List, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PPT = ROOT / "Retail_AI_Presentation.pptx"
COUNTS_JSON = ROOT / "outputs" / "counts_summary.json"
INPUT_IMAGE = ROOT / "cooldrinks.jpg"
DEFAULT_ANNOTATED = ROOT / "outputs" / "annotated_result.jpg"
OUTPUTS_DIR = ROOT / "outputs"


COLOR_BG = RGBColor(244, 248, 252)
COLOR_PRIMARY = RGBColor(16, 64, 104)
COLOR_SECONDARY = RGBColor(36, 87, 132)
COLOR_TEXT = RGBColor(33, 37, 41)
COLOR_ACCENT = RGBColor(0, 121, 140)


def set_slide_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG


def add_title(slide, title: str, subtitle: str = "") -> None:
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12.0), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.18), Inches(11.5), Inches(0.55))
        sf = sub_box.text_frame
        sf.clear()
        sp = sf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(16)
        sp.font.color.rgb = COLOR_SECONDARY


def add_bullets(slide, bullets: List[str], left: float = 0.9, top: float = 1.9, width: float = 11.5, height: float = 4.6) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(24)
        p.font.color.rgb = COLOR_TEXT


def add_footer(slide, text: str) -> None:
    fbox = slide.shapes.add_textbox(Inches(0.7), Inches(6.75), Inches(12.0), Inches(0.35))
    ft = fbox.text_frame
    ft.clear()
    p = ft.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(90, 98, 104)
    p.alignment = PP_ALIGN.RIGHT


def load_counts() -> Dict[str, int]:
    if not COUNTS_JSON.exists():
        return {}
    try:
        payload = json.loads(COUNTS_JSON.read_text(encoding="utf-8"))
    except Exception:
        return {}
    totals = payload.get("totals", {})
    if not isinstance(totals, dict):
        return {}
    result: Dict[str, int] = {}
    for k, v in totals.items():
        try:
            result[str(k)] = int(v)
        except Exception:
            continue
    return result


def resolve_annotated_image() -> Path | None:
    if DEFAULT_ANNOTATED.exists():
        return DEFAULT_ANNOTATED
    candidates = sorted(OUTPUTS_DIR.glob("annotated_*.jpg"), key=lambda p: p.stat().st_mtime, reverse=True)
    return candidates[0] if candidates else None


def add_pipeline_shapes(slide) -> None:
    steps = [
        "Input Image",
        "YOLO Detection",
        "Crop Regions",
        "CLIP Classification",
        "Count Aggregation",
        "Planogram + Inventory",
    ]
    left = 0.65
    top = 2.2
    w = 1.9
    h = 1.0
    gap = 0.15

    for i, step in enumerate(steps):
        x = left + i * (w + gap)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(top),
            Inches(w),
            Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(223, 237, 248) if i % 2 == 0 else RGBColor(203, 228, 243)
        shape.line.color.rgb = COLOR_SECONDARY
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = step
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY

        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                Inches(x + w),
                Inches(top + 0.35),
                Inches(gap),
                Inches(0.3),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_ACCENT
            arrow.line.color.rgb = COLOR_ACCENT


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1) Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    band = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.6))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(230, 241, 249)
    band.line.fill.background()
    add_title(s, "Shelf Product Detection and Inventory Intelligence", "Capstone Project Demo | Group 1 | Updated")
    add_bullets(
        s,
        [
            "YOLO + CLIP for Product Detection, Classification, and Counting",
            "Planogram-Aware Stock Estimation with Streamlit Demo Interface",
        ],
        top=2.25,
        height=2.0,
    )
    add_footer(s, "Academic Year 2026")

    # 2) Problem statement
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "Problem Statement", "Why this project matters")
    add_bullets(
        s,
        [
            "Manual shelf audits are slow, repetitive, and error-prone",
            "Stockouts happen because shelf visibility is not continuously tracked",
            "Retail teams need a fast way to convert shelf images into actionable inventory data",
            "Goal: Automate detection, counting, and reorder insights from shelf photos",
        ],
    )

    # 3) Objectives
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "Project Objectives")
    add_bullets(
        s,
        [
            "Detect products on shelf images using YOLO",
            "Classify each detected region with CLIP label prompts",
            "Generate product-wise counts and export CSV/JSON reports",
            "Support inventory updates and planogram-based quantity estimation",
            "Provide operator-friendly web interface for real-time usage",
        ],
    )

    # 4) Architecture
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "System Architecture", "End-to-end processing pipeline")
    add_pipeline_shapes(s)
    add_bullets(
        s,
        [
            "Auto profile selection chooses Beverage vs Shampoo labels",
            "Outputs include annotated images, count summaries, and reorder indicators",
        ],
        top=4.15,
        height=1.5,
    )

    # 5) Tech stack
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "Technology Stack")
    add_bullets(
        s,
        [
            "Language: Python 3.11",
            "Detection: Ultralytics YOLO (best.pt)",
            "Classification: OpenAI CLIP (clip-vit-base-patch32)",
            "Web App: Streamlit (Detection Studio + Inventory Center)",
            "Data and Reporting: Pandas, JSON, CSV, Pillow",
        ],
    )

    # 6) Methodology
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "Methodology")
    add_bullets(
        s,
        [
            "1. Detect product bounding boxes with confidence and NMS thresholds",
            "2. Crop each detection and run CLIP text-image similarity",
            "3. Apply corrections and non-product filtering prompts",
            "4. Aggregate counts by Category | Product",
            "5. Optionally sync counts into inventory records",
        ],
    )

    # 7) Planogram logic
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "Planogram-Based Quantity Estimation")
    add_bullets(
        s,
        [
            "Detected facings = d, Expected facings = e",
            "Compliance ratio r = min(1, d / e)",
            "Estimated units = round(r x full_stock_units)",
            "Reorder qty = max(0, reorder_level - estimated_units)",
            "Needs reorder if estimated_units < reorder_level",
        ],
    )

    # 8) Results table
    counts = load_counts()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "Sample Results", "Current run summary from project outputs")

    table_data: List[Tuple[str, int]] = sorted(counts.items(), key=lambda kv: kv[1], reverse=True)
    if not table_data:
        add_bullets(s, ["No counts_summary.json found. Run detection to populate results."])
    else:
        rows = min(1 + len(table_data), 7)
        cols = 3
        tbl = s.shapes.add_table(rows, cols, Inches(1.0), Inches(2.0), Inches(11.4), Inches(3.2)).table
        tbl.columns[0].width = Inches(5.0)
        tbl.columns[1].width = Inches(4.2)
        tbl.columns[2].width = Inches(2.2)

        tbl.cell(0, 0).text = "Category"
        tbl.cell(0, 1).text = "Product"
        tbl.cell(0, 2).text = "Count"

        for c in range(cols):
            cell = tbl.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(214, 233, 246)
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = COLOR_PRIMARY

        for i, (key, value) in enumerate(table_data[: rows - 1], start=1):
            if "|" in key:
                category, product = [p.strip() for p in key.split("|", 1)]
            else:
                category, product = "Unknown", key
            tbl.cell(i, 0).text = category
            tbl.cell(i, 1).text = product
            tbl.cell(i, 2).text = str(value)
            for c in range(cols):
                p = tbl.cell(i, c).text_frame.paragraphs[0]
                p.font.size = Pt(13)
                p.font.color.rgb = COLOR_TEXT

    # 9) Visual evidence
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "Visual Output", "Input image and annotated detection result")

    annotated = resolve_annotated_image()
    if INPUT_IMAGE.exists():
        s.shapes.add_picture(str(INPUT_IMAGE), Inches(0.8), Inches(1.9), width=Inches(5.8), height=Inches(4.5))
        label1 = s.shapes.add_textbox(Inches(2.8), Inches(6.45), Inches(1.8), Inches(0.3))
        p1 = label1.text_frame.paragraphs[0]
        p1.text = "Input"
        p1.alignment = PP_ALIGN.CENTER
        p1.font.bold = True
        p1.font.size = Pt(13)

    if annotated and annotated.exists():
        s.shapes.add_picture(str(annotated), Inches(6.8), Inches(1.9), width=Inches(5.8), height=Inches(4.5))
        label2 = s.shapes.add_textbox(Inches(8.9), Inches(6.45), Inches(2.0), Inches(0.3))
        p2 = label2.text_frame.paragraphs[0]
        p2.text = "Annotated Output"
        p2.alignment = PP_ALIGN.CENTER
        p2.font.bold = True
        p2.font.size = Pt(13)

    # 10) Streamlit app features
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "Live Demo Interface (Streamlit)")
    add_bullets(
        s,
        [
            "Detection Studio: upload/camera input, run inference, view counts",
            "Inventory Center: manage stock, reorder levels, and low-stock alerts",
            "Planogram Setup: expected facings, full-stock units, reorder thresholds",
            "One-click actions: reconcile stock from detection or planogram estimate",
        ],
    )

    # 11) Limitations and future work
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "Limitations and Future Enhancements")
    add_bullets(
        s,
        [
            "Current accuracy depends on image quality and product occlusions",
            "Similar packaging can still cause classification confusion",
            "Future: video tracking, stronger brand-specific fine-tuning",
            "Future: slot-level planogram constraints and trend dashboards",
        ],
    )

    # 12) Conclusion
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_title(s, "Conclusion")
    add_bullets(
        s,
        [
            "The project delivers an end-to-end retail shelf intelligence workflow",
            "YOLO + CLIP enables automated detection, counting, and reporting",
            "Planogram-aware estimation bridges computer vision with inventory decisions",
            "Thank You - Questions and Discussion",
        ],
    )

    prs.save(str(OUTPUT_PPT))
    return OUTPUT_PPT


if __name__ == "__main__":
    out = build_presentation()
    print(f"Generated presentation: {out}")
